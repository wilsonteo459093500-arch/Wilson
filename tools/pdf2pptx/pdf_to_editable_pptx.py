#!/usr/bin/env python3
"""
把设计稿 PDF 转成"可编辑"的 PPTX，且尽量 1:1 保留原版式。

思路（High-fidelity PDF -> PPTX）:
  1. 复制一份 PDF，用 redaction 只删除文字（图片 / 矢量图形原样保留），
     再按高倍率渲染成整页底图 -> 保证配图、色块、线条的分辨率与位置完全一致。
  2. 从原 PDF 精确读取每一段文字的基线坐标、字号、颜色、字距，
     在 PPT 上用真正的文本框重建 -> 文字可编辑、可搜索、可复制。
  3. 坐标按 (幻灯片宽 / 页面宽) 等比缩放，因此对齐关系与原稿完全一致。

用法:
    python3 pdf_to_editable_pptx.py input.pdf output.pptx [--dpi-scale 2.0]
"""

from __future__ import annotations

import argparse
import collections
import io
import json
import math
import statistics
import sys
from dataclasses import dataclass, field

import pymupdf
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
XML_NS = "http://www.w3.org/XML/1998/namespace"


def a(tag: str) -> str:
    return "{%s}%s" % (A_NS, tag)


# --------------------------------------------------------------------------
# 字体映射：PDF 内嵌字体名 -> PowerPoint 里使用的字体名
#   desc = 该字体的下伸部比例(em)，用于把"基线坐标"换算成"文本框顶部坐标"
# --------------------------------------------------------------------------
FONT_MAP = {
    "SourceHanSansCN-Normal": {"name": "Source Han Sans CN Normal", "bold": False, "desc": 0.288},
    "SourceHanSansCN-Light": {"name": "Source Han Sans CN Light", "bold": False, "desc": 0.288},
    "SourceHanSansCN-Regular": {"name": "Source Han Sans CN", "bold": False, "desc": 0.288},
    "SourceHanSerifCN-Regular": {"name": "Source Han Serif CN", "bold": False, "desc": 0.286},
    "PingFangSC-Ultralight-B5pc-H": {"name": "PingFang SC Ultralight", "bold": False, "desc": 0.24},
    "youyou-yisong": {"name": "youyou-yisong", "bold": False, "desc": 0.20},
    "SabonLTStd-Bold": {"name": "Sabon LT Std", "bold": True, "desc": 0.228},
    "SabonLTStd-Roman": {"name": "Sabon LT Std", "bold": False, "desc": 0.228},
}
DEFAULT_DESC = 0.24

# PDF 里数字被映射到了私用区(旧式数字字形)，还原成普通数字，方便编辑
PUA_DIGITS = {0xF6B1 + i: str(i) for i in range(10)}
LIGATURES = {"ﬀ": "ff", "ﬁ": "fi", "ﬂ": "fl", "ﬃ": "ffi", "ﬄ": "ffl", "ﬅ": "ft", "ﬆ": "st"}


def clean_text(s: str) -> str:
    out = []
    for ch in s:
        if ord(ch) in PUA_DIGITS:
            out.append(PUA_DIGITS[ord(ch)])
        elif ch in LIGATURES:
            out.append(LIGATURES[ch])
        else:
            out.append(ch)
    return "".join(out)


def load_font_map(path: str | None):
    """允许用 JSON 覆盖字体映射: {"youyou-yisong": {"name": "悠悠宋", "bold": false, "desc": 0.2}}"""
    if not path:
        return
    with open(path, encoding="utf-8") as fh:
        for k, v in json.load(fh).items():
            base = FONT_MAP.get(k, {"bold": False, "desc": DEFAULT_DESC}).copy()
            base.update(v)
            FONT_MAP[k] = base


def font_info(pdf_font: str) -> dict:
    base = pdf_font.split("+")[-1]
    if base in FONT_MAP:
        return FONT_MAP[base]
    guess = base.split("-")[0]
    return {"name": guess or base, "bold": "Bold" in base, "desc": DEFAULT_DESC}


# --------------------------------------------------------------------------
# 数据结构
# --------------------------------------------------------------------------
@dataclass
class Run:
    text: str
    font: str          # PDF 字体名
    size: float        # pt (PDF 坐标系)
    color: int
    track: float       # 字距, pt (PDF 坐标系)
    italic: bool


@dataclass
class Line:
    baseline: float
    x0: float
    x1: float
    size: float
    runs: list = field(default_factory=list)

    @property
    def center(self) -> float:
        return (self.x0 + self.x1) / 2.0

    @property
    def sig(self):
        r = self.runs[0]
        return (r.font, round(self.size, 1), r.color)


@dataclass
class Block:
    lines: list
    lead: float | None = None

    @property
    def size(self) -> float:
        return self.lines[0].size


# --------------------------------------------------------------------------
# 文字抽取
# --------------------------------------------------------------------------
def extract_lines(page, fontbufs) -> tuple[list, list]:
    """返回 (可编辑的行, 需要留在底图上的 span 矩形)"""
    lines, keep_baked = [], []
    for block in page.get_text("rawdict")["blocks"]:
        if block["type"] != 0:
            continue
        for ln in block["lines"]:
            horizontal = abs(ln["dir"][0] - 1.0) < 1e-6 and abs(ln["dir"][1]) < 1e-6
            runs, x0, x1, baseline, size = [], None, None, None, 0.0
            for sp in ln["spans"]:
                text = clean_text("".join(c["c"] for c in sp["chars"]))
                if not text.strip():
                    # 纯空白 span：并入上一段，避免产生空文本框
                    if runs:
                        runs[-1].text += text
                    continue
                if not horizontal:
                    keep_baked.append(sp)
                    continue
                track = measure_tracking(sp, fontbufs)
                runs.append(
                    Run(
                        text=text,
                        font=sp["font"],
                        size=sp["size"],
                        color=sp["color"],
                        track=track,
                        italic=bool(sp["flags"] & 2),
                    )
                )
                sx0 = sp["chars"][0]["origin"][0]
                x0 = sx0 if x0 is None else min(x0, sx0)
                x1 = sp["bbox"][2] if x1 is None else max(x1, sp["bbox"][2])
                baseline = sp["chars"][0]["origin"][1] if baseline is None else baseline
                size = max(size, sp["size"])
            if runs:
                lines.append(Line(baseline=baseline, x0=x0, x1=x1, size=size, runs=runs))
    lines.sort(key=lambda l: (round(l.baseline, 1), l.x0))
    return lines, keep_baked


def measure_tracking(span, fontbufs) -> float:
    """从字符实际步进与字体自然步进之差，反推该 span 的字距(pt)。"""
    fb = fontbufs.get(span["font"].split("+")[-1])
    chars = span["chars"]
    if fb is None or len(chars) < 3:
        return 0.0
    diffs = []
    for i in range(len(chars) - 1):
        if chars[i]["c"] == " " or chars[i + 1]["c"] == " ":
            continue
        adv = chars[i + 1]["origin"][0] - chars[i]["origin"][0]
        if adv <= 0:
            continue
        nat = fb.text_length(chars[i]["c"], fontsize=span["size"])
        diffs.append(adv - nat)
    if len(diffs) < 2:
        return 0.0
    med = statistics.median(diffs)
    return med if abs(med) > 0.02 else 0.0


# --------------------------------------------------------------------------
# 把行合并成段落块（同字体/同字号/同颜色/等行距/同对齐）
# --------------------------------------------------------------------------
def group_lines(lines: list) -> list:
    blocks: list[Block] = []
    for ln in lines:
        placed = False
        for blk in reversed(blocks):
            last = blk.lines[-1]
            if last.sig != ln.sig:
                continue
            delta = ln.baseline - last.baseline
            if delta <= 0.5 or delta > 2.7 * ln.size:
                continue
            if blk.lead is not None and abs(delta - blk.lead) > 0.6:
                continue
            if not (
                abs(ln.x0 - blk.lines[0].x0) <= 0.9
                or abs(ln.center - blk.lines[0].center) <= 0.9
                or abs(ln.x1 - blk.lines[0].x1) <= 0.9
            ):
                continue
            if blk.lead is None:
                blk.lead = delta
            blk.lines.append(ln)
            placed = True
            break
        if not placed:
            blocks.append(Block(lines=[ln]))
    return blocks


def block_align(blk: Block, page_width: float) -> str:
    if len(blk.lines) >= 2:
        xs = [l.x0 for l in blk.lines]
        cs = [l.center for l in blk.lines]
        rs = [l.x1 for l in blk.lines]
        if max(xs) - min(xs) <= 0.9:
            return "l"
        if max(cs) - min(cs) <= 0.9:
            return "ctr"
        if max(rs) - min(rs) <= 0.9:
            return "r"
        return "l"
    # 单行：若相对整页居中，保持居中，编辑后依然居中
    if abs(blk.lines[0].center - page_width / 2.0) <= 4.0:
        return "ctr"
    return "l"


# --------------------------------------------------------------------------
# 底图：删掉文字后整页高倍渲染
# --------------------------------------------------------------------------
def render_background(src_path: str, pno: int, target_px: int, jpeg_quality: int):
    doc = pymupdf.open(src_path)
    page = doc[pno]
    for block in page.get_text("dict")["blocks"]:
        if block["type"] != 0:
            continue
        for ln in block["lines"]:
            if not (abs(ln["dir"][0] - 1.0) < 1e-6 and abs(ln["dir"][1]) < 1e-6):
                continue  # 竖排文字保留在底图上
            for sp in ln["spans"]:
                r = pymupdf.Rect(sp["bbox"])
                r.x0 -= 1.2
                r.x1 += 1.2
                r.y0 -= 2.0
                r.y1 += 2.0
                page.add_redact_annot(r, fill=False)
    page.apply_redactions(
        images=pymupdf.PDF_REDACT_IMAGE_NONE,
        graphics=pymupdf.PDF_REDACT_LINE_ART_NONE,
        text=pymupdf.PDF_REDACT_TEXT_REMOVE,
    )
    zoom = target_px / page.rect.width
    pix = page.get_pixmap(matrix=pymupdf.Matrix(zoom, zoom), alpha=False)
    png = pix.tobytes("png")
    img = Image.open(io.BytesIO(png)).convert("RGB")
    jpg_buf = io.BytesIO()
    img.save(jpg_buf, "JPEG", quality=jpeg_quality, subsampling=0, optimize=True)
    doc.close()
    jpg = jpg_buf.getvalue()
    if len(png) <= len(jpg):
        return png, "png"
    return jpg, "jpg"


# --------------------------------------------------------------------------
# PPTX 写入
# --------------------------------------------------------------------------
def set_lock(pic):
    cNvPicPr = pic._element.nvPicPr.find(a("cNvPicPr"))
    if cNvPicPr is None:
        return
    locks = etree.SubElement(cNvPicPr, a("picLocks"))
    locks.set("noMove", "1")
    locks.set("noResize", "1")
    locks.set("noChangeAspect", "1")


def add_textbox(slide, blk: Block, scale: float, page_width: float, off_x: float, off_y: float):
    info0 = font_info(blk.lines[0].runs[0].font)
    size0 = blk.size
    lead = blk.lead if blk.lead else max(size0 * 1.2, size0 + 1.0)
    align = block_align(blk, page_width)

    # PowerPoint(固定行距)下: 首行基线 = 文本框顶 + 行距 - 下伸部
    top_pdf = blk.lines[0].baseline - lead + info0["desc"] * size0
    width_pdf = max(l.x1 - l.x0 for l in blk.lines)
    height_pdf = lead * len(blk.lines) + size0 * 0.6

    # 预留一点编辑空间（关闭自动换行，加字不会重排，只是点选区域大一点）
    pad = max(width_pdf * 0.12, size0 * 1.5)
    if align == "l":
        left_pdf = blk.lines[0].x0
        box_w = width_pdf + pad
    elif align == "ctr":
        center = statistics.mean([l.center for l in blk.lines])
        if len(blk.lines) == 1 and abs(center - page_width / 2.0) <= 4.0:
            center = page_width / 2.0
        box_w = width_pdf + pad
        left_pdf = center - box_w / 2.0
    else:
        right = max(l.x1 for l in blk.lines)
        box_w = width_pdf + pad
        left_pdf = right - box_w

    box = slide.shapes.add_textbox(
        Emu(int(round((off_x + left_pdf * scale) * 914400))),
        Emu(int(round((off_y + top_pdf * scale) * 914400))),
        Emu(int(round(box_w * scale * 914400))),
        Emu(int(round(height_pdf * scale * 914400))),
    )
    tf = box.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    body = tf._txBody.find(a("bodyPr"))
    body.set("anchor", "t")
    body.set("anchorCtr", "0")

    for i, ln in enumerate(blk.lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pPr = para._p.get_or_add_pPr()
        pPr.set("marL", "0")
        pPr.set("indent", "0")
        pPr.set("algn", align)
        lnSpc = etree.SubElement(pPr, a("lnSpc"))
        spcPts = etree.SubElement(lnSpc, a("spcPts"))
        spcPts.set("val", str(int(round(lead * scale * 72 * 100))))
        for key in ("spcBef", "spcAft"):
            el = etree.SubElement(pPr, a(key))
            pts = etree.SubElement(el, a("spcPts"))
            pts.set("val", "0")
        for r in ln.runs:
            run = para.add_run()
            run.text = r.text
            t = run._r.find(a("t"))
            if t is not None and (r.text != r.text.strip()):
                t.set("{%s}space" % XML_NS, "preserve")
            fi = font_info(r.font)
            f = run.font
            f.size = Pt(round(r.size * scale * 72, 2))
            f.bold = fi["bold"]
            f.italic = r.italic or None
            f.color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor(
                (r.color >> 16) & 0xFF, (r.color >> 8) & 0xFF, r.color & 0xFF
            )
            f.name = fi["name"]
            rPr = run._r.get_or_add_rPr()
            track_pt = r.track * scale * 72
            if abs(track_pt) >= 0.01:
                rPr.set("spc", str(int(round(track_pt * 100))))
            latin = rPr.find(a("latin"))
            for tag in ("ea", "cs"):
                el = etree.SubElement(rPr, a(tag))
                el.set("typeface", fi["name"])
                if latin is not None:
                    latin.addnext(el)
        # 让"在段尾继续输入"时保持同样的字体样式
        last_rPr = para._p.findall(a("r"))[-1].find(a("rPr"))
        end = etree.SubElement(para._p, a("endParaRPr"))
        for k, v in last_rPr.attrib.items():
            end.set(k, v)
        for child in last_rPr:
            end.append(etree.fromstring(etree.tostring(child)))
    return box


def convert(src: str, dst: str, target_px: int, jpeg_quality: int, report_path: str | None):
    doc = pymupdf.open(src)

    fontbufs = {}
    for pno in range(doc.page_count):
        for f in doc[pno].get_fonts(full=True):
            base = f[3].split("+")[-1]
            if base in fontbufs:
                continue
            try:
                _, _, _, buf = doc.extract_font(f[0])
                if buf:
                    fontbufs[base] = pymupdf.Font(fontbuffer=buf)
            except Exception:
                pass

    page0 = doc[0]
    slide_w_in = 13.3333
    slide_h_in = round(slide_w_in * page0.rect.height / page0.rect.width, 4)
    prs = Presentation()
    prs.slide_width = Emu(int(round(slide_w_in * 914400)))
    prs.slide_height = Emu(int(round(slide_h_in * 914400)))
    blank = prs.slide_layouts[6]

    report = {
        "source": src,
        "slide_size_in": [slide_w_in, slide_h_in],
        "background_px": target_px,
        "pages": [],
        "fonts": sorted({font_info(k)["name"] for k in fontbufs}),
    }

    for pno in range(doc.page_count):
        page = doc[pno]
        pw, ph = page.rect.width, page.rect.height
        scale = slide_w_in / pw                    # inch per pdf-unit
        off_x = 0.0
        off_y = (slide_h_in - ph * scale) / 2.0    # 页面比例略有差异时垂直居中

        slide = prs.slides.add_slide(blank)

        img_bytes, ext = render_background(src, pno, target_px, jpeg_quality)
        pic = slide.shapes.add_picture(
            io.BytesIO(img_bytes),
            Emu(int(round(off_x * 914400))),
            Emu(int(round(off_y * 914400))),
            Emu(int(round(pw * scale * 914400))),
            Emu(int(round(ph * scale * 914400))),
        )
        pic.name = f"Background p{pno + 1}"
        set_lock(pic)

        lines, baked = extract_lines(page, fontbufs)
        blocks = group_lines(lines)
        for blk in blocks:
            add_textbox(slide, blk, scale, pw, off_x, off_y)

        report["pages"].append(
            {
                "page": pno + 1,
                "text_boxes": len(blocks),
                "lines": len(lines),
                "baked_spans": len(baked),
                "bg_format": ext,
                "bg_kb": round(len(img_bytes) / 1024),
            }
        )
        print(
            f"  p{pno + 1:>2}: {len(blocks):>3} text boxes / {len(lines):>3} lines, "
            f"bg {ext} {len(img_bytes) // 1024}KB",
            flush=True,
        )

    prs.save(dst)
    if report_path:
        json.dump(report, open(report_path, "w"), ensure_ascii=False, indent=2)
    return report


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("src")
    ap.add_argument("dst")
    ap.add_argument("--bg-width", type=int, default=3840, help="底图像素宽度")
    ap.add_argument("--jpeg-quality", type=int, default=90)
    ap.add_argument("--report", default=None)
    ap.add_argument("--font-map", default=None, help="JSON: PDF 字体名 -> PPT 字体名")
    args = ap.parse_args()
    load_font_map(args.font_map)
    print(f"converting {args.src} -> {args.dst}")
    convert(args.src, args.dst, args.bg_width, args.jpeg_quality, args.report)
    print("done")


if __name__ == "__main__":
    main()
